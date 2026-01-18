import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import docx
import openpyxl
import re
import io
import dashscope
import json
import time

# ==========================================
# 0. 全局配置与初始化
# ==========================================
st.set_page_config(page_title="行研 Copilot Pro", layout="wide", page_icon="🚀")

# 设置 Matplotlib 中文字体 (防止方块乱码)
plt.rcParams['font.sans-serif'] = ['SimHei', 'Arial', 'Microsoft YaHei', 'PingFang SC']
plt.rcParams['axes.unicode_minus'] = False

# 初始化 Session State (用于跨刷新保存数据)
if 'ai_config' not in st.session_state:
    st.session_state['ai_config'] = None
if 'df_cache' not in st.session_state:
    st.session_state['df_cache'] = None

# ==========================================
# 1. 基础工具函数：全格式文档解析
# ==========================================
def clean_text(text):
    """清洗文本：去除空格和换行"""
    if not text: return ""
    return "".join(str(text).split())

def split_segments(full_text):
    """将长文本切分为短句集合，用于比对"""
    segments = re.split(r'[。；！？\n]+', str(full_text))
    return set([clean_text(s) for s in segments if len(clean_text(s)) > 2])

def get_docx_text(file):
    """解析 Word"""
    try:
        doc = docx.Document(file)
        txt = []
        for p in doc.paragraphs: txt.append(p.text)
        for t in doc.tables:
            for r in t.rows:
                for c in r.cells: txt.append(c.text)
        raw = "\n".join(txt)
        return split_segments(raw), raw
    except: return set(), ""

def get_pptx_text(file):
    """解析 PPT"""
    try:
        prs = Presentation(file)
        txt = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"): txt.append(shape.text)
                if shape.has_table:
                    for r in shape.table.rows:
                        for c in r.cells: txt.append(c.text)
        raw = "\n".join(txt)
        return split_segments(raw), raw
    except: return set(), ""

def get_excel_text(file):
    """解析 Excel"""
    try:
        wb = openpyxl.load_workbook(file, data_only=True)
        txt = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                for cell in row:
                    if cell: txt.append(str(cell))
        raw = "\n".join(txt)
        return split_segments(raw), raw
    except: return set(), ""
    
@st.cache_data(show_spinner=False)
def dispatch_extractor(file):
    """根据后缀分发解析器"""
    if file.name.endswith('.docx'): return get_docx_text(file)
    elif file.name.endswith('.pptx'): return get_pptx_text(file)
    elif file.name.endswith('.xlsx'): return get_excel_text(file)
    return set(), ""

# ==========================================
# 2. 模块：全格式文档核对 (Compliance Check)
# ==========================================
def module_compliance(api_key):
    st.header("🕵️ 全格式文档核对")
    st.markdown("支持 **Word / Excel / PPT** 任意格式两两比对。")
    st.caption("场景：检查 PPT 里的数据是否和 Excel 底稿一致，或检查 PPT 是否有 Word 报告里没有的“鬼话”。")
    
    c1, c2 = st.columns(2)
    # 添加 unique key 防止报错
    f1 = c1.file_uploader("1. 基准文件 (Source)", type=['docx','xlsx','pptx'], key="comp_f1")
    f2 = c2.file_uploader("2. 待测文件 (Target)", type=['docx','xlsx','pptx'], key="comp_f2")
    
    if f1 and f2:
        with st.spinner("正在抽取文本并比对..."):
            set1, raw1 = dispatch_extractor(f1)
            set2, raw2 = dispatch_extractor(f2)
            
            # 核心逻辑：Target - Source = 鬼话
            ghosts = list(set2 - set1)
            
        if not ghosts:
            st.success("✅ 完美匹配！Target 中的所有内容均可在 Source 中找到依据。")
        else:
            st.warning(f"⚠️ 发现 {len(ghosts)} 处内容来源不明")
            
            # AI 仲裁
            if st.button("🧠 AI 语义判别 (是改写还是瞎编?)", key="btn_comp_ai"):
                if not api_key:
                    st.error("请先在左侧侧边栏输入 API Key")
                else:
                    with st.spinner("AI 正在阅读文档..."):
                        dashscope.api_key = api_key
                        # 截取原文前 5000 字防止超长
                        prompt = f"""
                        你是一个金融合规审核员。
                        【基准事实】：{raw1[:5000]}...
                        【待审核内容】：{ghosts[:15]}
                        
                        任务：判断待审核内容是否在基准事实中有依据？
                        1. 如果是同义改写或数据一致，请标记为【通过】。
                        2. 如果完全不存在或数据冲突，请标记为【存疑】。
                        请直接输出分析结果。
                        """
                        try:
                            resp = dashscope.Generation.call(model='qwen-turbo', prompt=prompt)
                            st.info(resp.output.text)
                        except Exception as e:
                            st.error(f"AI 调用失败: {e}")
            
            with st.expander("🔍 查看详细差异列表"):
                st.write(ghosts)

# ==========================================
# 3. 模块：智能会议纪要 (Q&A版)
# ==========================================
def module_meeting(api_key):
    st.header("🎙️ 智能会议纪要 (Q&A结构化)")
    st.caption("上传录音 -> 自动转写 -> 生成【核心观点】+【问答实录】。")
    
    audio_file = st.file_uploader("上传录音文件", type=['mp3','wav','m4a'], key="meet_audio")
    
    if audio_file and st.button("开始分析", key="btn_meet_gen"):
        if not api_key:
            st.error("需要 API Key 才能使用 AI 功能")
            return
            
        st.info("🔄 正在进行语音识别 (ASR)...")
        # 模拟等待效果
        time.sleep(1.5)
        
        # --- 模拟识别结果 (真实场景需对接阿里云 OSS + ASR 接口) ---
        mock_text = """
        王总：大家好。Q1我们营收100亿，同比增长20%，净利润15亿。下面开始问答。
        分析师张三：请问毛利率为什么下降了？
        王总：主要是因为原材料铜价上涨了15%，压缩了利润空间，但我们通过套保对冲了一部分。
        分析师李四：未来产能规划如何？
        王总：我们预计下半年随着越南新产线投产，产能将提升30%，毛利率会回升到30%。
        """
        # ---------------------------------------------------
        
        st.success("✅ 语音识别完成！")
        with st.expander("查看识别原文"):
            st.text(mock_text)
        
        st.info("🧠 AI 正在整理 Q&A 结构...")
        dashscope.api_key = api_key
        prompt = f"""
        你是一个行研分析师。请将以下会议文本整理为规范的会议纪要。
        
        【要求】
        1. 总结核心要点 (Bullet points)。
        2. Q&A 环节必须严格按照 "Q: [问题] \n A: [回答]" 的格式整理。
        3. 去除口语废话，保持专业性。
        
        【会议文本】：{mock_text}
        """
        try:
            resp = dashscope.Generation.call(model='qwen-turbo', prompt=prompt)
            st.markdown("### 📝 纪要预览")
            st.markdown(resp.output.text)
            st.download_button("📥 下载纪要 TXT", resp.output.text, "minutes.txt", key="btn_dl_txt")
        except Exception as e:
            st.error(f"生成失败: {e}")

# ==========================================
# 4. 模块：智能制图 (范例仿制版)
# ==========================================
@st.cache_data(show_spinner=False)
def ai_analyze_chart(api_key, df):
    """AI 分析引擎：决定怎么画"""
    dashscope.api_key = api_key
    data_sample = df.head(3).to_json(orient='records', force_ascii=False)
    prompt = f"""
    分析以下 Excel 数据样例，给出 Matplotlib 绘图建议。
    数据：{data_sample}
    请严格返回如下 JSON 格式 (不要 Markdown)：
    {{
        "chart_type": "dual_axis" 或 "bar" 或 "line",
        "x_col": "推测的时间或类别列名",
        "y_primary": ["主轴列名1"],
        "y_secondary": ["副轴列名1"] (仅双轴图需要, 否则为空list),
        "title": "建议标题",
        "summary": "一句话数据洞察"
    }}
    """
    try:
        resp = dashscope.Generation.call(model='qwen-turbo', prompt=prompt)
        txt = resp.output.text.replace("```json","").replace("```","").strip()
        return json.loads(txt)
    except: return None

def module_smart_chart_ref(api_key):
    st.header("📊 智能制图 (范例仿制版)")
    st.markdown("**工作流：** 上传参考截图 -> 上传数据 -> 调整颜色以匹配参考图 -> 导出。")
    
    c1, c2 = st.columns(2)
    # 修改：上传参考图片而不是PPT模板
    ref_image = c1.file_uploader("1. 上传参考范例 (截图)", type=['png', 'jpg', 'jpeg'], key="chart_ref_img")
    data_file = c2.file_uploader("2. 上传新数据 (Excel)", type=['xlsx'], key="chart_data_excel")
    
    # 辅助显示参考图
    if ref_image:
        with c1.expander("👁️ 参考图预览 (对照调整下方颜色)", expanded=True):
            st.image(ref_image, use_column_width=True)

    # 数据加载与 AI 分析
    if data_file and api_key:
        df = pd.read_excel(data_file)
        st.session_state['df_cache'] = df
        
        # 按钮：触发 AI
        if st.button("🤖 1. AI 分析数据结构", type="primary", key="btn_ai_analyze"):
            with st.spinner("AI 正在思考最佳画法..."):
                config = ai_analyze_chart(api_key, df)
                if config:
                    st.session_state['ai_config'] = config
                    st.success("分析完成！请在下方进行【样式对齐】。")
                else:
                    st.error("AI 分析失败，请检查 API Key")

    # 微调面板
    if st.session_state['ai_config']:
        config = st.session_state['ai_config']
        df = st.session_state['df_cache']
        cols = df.columns.tolist()

        st.divider()
        st.subheader("🎨 2. 样式对齐 (Style Alignment)")
        
        # 布局：左控右图
        col_ctrl, col_view = st.columns([1, 2])
        
        with col_ctrl:
            st.markdown("#### ⚙️ 图表参数")
            user_chart_type = st.selectbox("图表类型", ["dual_axis", "bar", "line"], 
                                           index=["dual_axis", "bar", "line"].index(config.get('chart_type', 'bar')),
                                           key="sel_chart_type")
            
            user_x = st.selectbox("X轴数据", cols, index=cols.index(config.get('x_col')) if config.get('x_col') in cols else 0, key="sel_x")
            
            default_y1 = [c for c in config.get('y_primary', []) if c in cols]
            user_y1 = st.multiselect("左轴/主数据", cols, default=default_y1 if default_y1 else [cols[1]], key="sel_y1")
            
            user_y2 = []
            if user_chart_type == "dual_axis":
                default_y2 = [c for c in config.get('y_secondary', []) if c in cols]
                user_y2 = st.multiselect("右轴/次数据", cols, default=default_y2, key="sel_y2")

            st.markdown("---")
            st.markdown("#### 🎨 风格复刻")
            st.caption("请吸取参考图的颜色填入下方：")
            
            color_1 = st.color_picker("主色调 (Bar/Left)", "#C00000", key="cp_1") 
            color_2 = st.color_picker("副色调 (Line/Right)", "#FFC000", key="cp_2")
            font_size = st.slider("字体大小", 8, 24, 12, key="sl_font")
            user_title = st.text_input("图表标题", value=config.get('title', 'Chart Title'), key="txt_title")

        # 实时绘图逻辑
        with col_view:
            st.markdown("#### 🖼️ 效果预览")
            plt.rcParams.update({'font.size': font_size})
            fig, ax1 = plt.subplots(figsize=(8, 4.5)) # 16:9 比例
            
            # 绘图核心逻辑
            if user_chart_type == "dual_axis":
                ax1.bar(df[user_x], df[user_y1[0]], color=color_1, alpha=0.8, label=user_y1[0])
                ax1.set_ylabel(user_y1[0], color=color_1, fontweight='bold')
                if user_y2:
                    ax2 = ax1.twinx()
                    ax2.plot(df[user_x], df[user_y2[0]], color=color_2, marker='o', linewidth=3, label=user_y2[0])
                    ax2.set_ylabel(user_y2[0], color=color_2, fontweight='bold')
                    ax2.grid(False)
            elif user_chart_type == "bar":
                for i, c in enumerate(user_y1):
                    ax1.bar(df[user_x], df[c], color=color_1 if i==0 else None, alpha=0.8, label=c)
            elif user_chart_type == "line":
                for i, c in enumerate(user_y1):
                    ax1.plot(df[user_x], df[c], color=color_2 if i==0 else None, marker='o', linewidth=2, label=c)

            ax1.set_title(user_title, pad=15, fontweight='bold')
            ax1.grid(True, linestyle='--', alpha=0.5)
            plt.tight_layout()
            
            st.pyplot(fig)
            
            # 保存逻辑
            img_stream = io.BytesIO()
            plt.savefig(img_stream, format='png', dpi=300, bbox_inches='tight')
            img_stream.seek(0)
            
            ppt_stream = io.BytesIO()
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6]) # 空白版式
            img_stream.seek(0)
            slide.shapes.add_picture(img_stream, Inches(0.5), Inches(0.5), width=Inches(9))
            prs.save(ppt_stream)
            ppt_stream.seek(0)
            img_stream.seek(0)

            # 下载区
            st.success("✅ 生成完毕！")
            d1, d2 = st.columns(2)
            d1.download_button("📥 下载高清 PNG", img_stream, "chart.png", "image/png", key="btn_dl_img")
            d2.download_button("📥 下载 PPT (含图)", ppt_stream, "chart.pptx", key="btn_dl_ppt")

# ==========================================
# 5. 主程序入口与导航
# ==========================================
with st.sidebar:
    st.title("🚀 行研 Copilot")
    st.caption("实习生的一站式工作台")
    api_key = st.text_input("🔑 API Key (Qwen)", type="password", key="main_api_key")
    st.markdown("---")
    mode = st.radio("功能导航", [
        "📊 智能制图 (范例仿制)",
        "🕵️ 全格式核对",
        "🎙️ 智能会议纪要"
    ], key="nav_radio")

if mode == "📊 智能制图 (范例仿制)":
    module_smart_chart_ref(api_key)
elif mode == "🕵️ 全格式核对":
    module_compliance(api_key)
elif mode == "🎙️ 智能会议纪要":
    module_meeting(api_key)